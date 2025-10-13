import csv
import os
import json
import datetime
from pptx import Presentation
from openai import OpenAI
import gspread
from google.oauth2.service_account import Credentials
from jinja2 import Environment, FileSystemLoader
# from weasyprint import HTML
from playwright.sync_api import sync_playwright


# ---------------------------
# CONFIG
# ---------------------------
TEACHER_SUBJECT = os.getenv("TEACHER_SUBJECT", "French")
PARCOURS_COUNT = os.getenv("PARCOURS_COUNT", 1)

CSV_FILE = "cahier_journal_2parcours.csv" if PARCOURS_COUNT == 2 else "cahier_journal.csv"
PPTX_FILE = "Français_Niv5_Parcour1_Palier3_Séance1.pptx"

OUTPUT_DIR = os.getenv("OUTPUT_DIR", "outputs_mindmaps")
PPTX_DIR = "./lessons"
SHEET_ID = os.getenv("SHEET_ID")  # from Google Sheets URL
os.makedirs(OUTPUT_DIR, exist_ok=True)

HEADERS_PARCOURS1 = [
    "Date",
    "Parcours",
    "Palier",
    "Séance",
    "Objectif",
    "Rituel",
    "Vocabulaire",
    "Lecture",
    "Écriture",
    "Pratique autonome",
    "Jeu",
    "Notes"
]
HEADERS_PARCOURS2 = [
    "Date",
    "Palier1", "Séance1", "Objectif1", "Rituel1", "Vocabulaire1", "Lecture1", "Écriture1", "Pratique autonome1", "Jeu1", "Notes1",
    "Palier2", "Séance2", "Objectif2", "Rituel2", "Vocabulaire2", "Lecture2", "Écriture2", "Pratique autonome2", "Jeu2", "Notes2",
]
HEADERS = HEADERS_PARCOURS2 if PARCOURS_COUNT == 2 else HEADERS_PARCOURS1
sample_csv_row = {
    "Date": "2025-10-05",
    "Parcours": "Parcours 1",
    "Palier": "3",
    "Séance": "1",
    "Objectif": "Amener les élèves à lire et écrire des mots simples contenant les lettres m, n, s, l.",
    "Rituel": "Chanson d'accueil, dictée flash de lettres et correction collective.",
    "Vocabulaire": "Découverte des mots : maman, lune, sel, sol, melon, salle.",
    "Lecture": "Lecture de syllabes et de mots simples avec m, n, s, l.",
    "Écriture": "Écriture guidée de mots et phrases au tableau puis sur ardoise.",
    "Pratique autonome": "Exercices individuels page 10 et lecture à voix basse par binômes.",
    "Jeu": "Jeu du mot mystère : deviner le mot à partir d’indices sonores.",
    "Notes": "La moitié des élèves confondent n/m, besoin de rebrassage demain."
}

csv_data_2parcours = {
    "Date": "2025-10-10",
    # Parcours 1
    "Palier1": "1",
    "Séance1": "3",
    "Objectif1": "Lire et écrire des mots avec la lettre m.",
    "Rituel1": "Dictée de syllabes et mots simples : ma, mi, mu.",
    "Vocabulaire1": "maman, maison, main.",
    "Lecture1": "Texte court sur le mot 'maman'.",
    "Écriture1": "Copie des mots appris et une phrase simple.",
    "Pratique autonome1": "Jeux de lecture sur ardoise.",
    "Jeu1": "Cherche le mot dans le texte.",
    "Notes1": "Bon engagement du groupe 1.",
    # Parcours 2
    "Palier2": "2",
    "Séance2": "3",
    "Objectif2": "Lire et écrire des phrases simples contenant la lettre j.",
    "Rituel2": "Dictée de syllabes et mots : ja, je, jou.",
    "Vocabulaire2": "jupe, jardin, jaune.",
    "Lecture2": "Lecture d’un petit texte : 'Le jardin de Julie'.",
    "Écriture2": "Écriture de 2 phrases avec les mots du jour.",
    "Pratique autonome2": "Jeu de cartes syllabes à associer.",
    "Jeu2": "Mots cachés avec la lettre j.",
    "Notes2": "Bon progrès en lecture, attention à l’écriture."
}


client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))


# ---------------------------
# HELPERS
# ---------------------------
def get_pptx_pairs(pptx_dir):
    parcours1 = sorted([f for f in os.listdir(pptx_dir) if "Parcour1" in f])
    parcours2 = sorted([f for f in os.listdir(pptx_dir) if "Parcour2" in f])
    pairs = list(zip(parcours1, parcours2))  # (p1, p2)
    return pairs

def ensure_csv_exists():
    """Check if CSV exists, if not create it with headers."""
    print("Checking for CSV file...")
    if not os.path.exists(CSV_FILE):
        with open(CSV_FILE, "w", newline="", encoding="utf-8") as f:
            print("Creating CSV file...")
            writer = csv.writer(f)
            writer.writerow(HEADERS)

def extract_title_info(pptx_path: str):
    """Extract the title from filename and palier/séance numbers."""
    print("Extracting title info...")
    filename = os.path.basename(pptx_path)
    name, _ = os.path.splitext(filename)

    palier, seance = None, None
    parts = name.split("_")
    for part in parts:
        if part.lower().startswith("palier"):
            palier = part.replace("Palier", "").replace("palier", "")
        if part.lower().startswith("séance") or part.lower().startswith("seance"):
            seance = part.replace("Séance", "").replace("séance", "").replace("Seance", "").replace("seance", "")

    title = name.replace("_", " ")
    print(f"Titre: {title} | Palier: {palier} | Séance: {seance if seance else '...'}")
    return title, palier, seance

def entry_exists_in_sheet(sheet, palier, seance):
    """Check if (Palier, Séance) already exist in the Google Sheet."""
    data = sheet.get_all_records()  # returns list[dict]
    for row in data:
        if str(row.get("Palier", "")).strip() == str(palier).strip() and str(row.get("Séance", "")).strip() == str(seance).strip():
            return True
    return False

def entry_exists_in_csv(palier, seance):
    """Check if (Palier, Séance) already exist in local CSV."""
    if not os.path.exists(CSV_FILE):
        return False
    with open(CSV_FILE, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            if str(row.get("Palier", "")).strip() == str(palier).strip() and str(row.get("Séance", "")).strip() == str(seance).strip():
                return True
    return False

def should_process(pptx_path: str) -> bool:
    """Return False if (Palier, Séance) already exist in CSV or Google Sheet."""
    title, palier, seance = extract_title_info(pptx_path)

    # If missing numbers, skip automatically
    if not palier or not seance:
        print(f"⚠️ Skipping {pptx_path} — missing Palier or Séance info.")
        return False

    # Check CSV
    if entry_exists_in_csv(palier, seance):
        print(f"⏩ Skipping {pptx_path} — already in CSV (Palier {palier}, Séance {seance}).")
        return False

    # Check Google Sheet (optional if you trust CSV as local cache)
    # SCOPES = ["https://www.googleapis.com/auth/spreadsheets"]
    # creds = Credentials.from_service_account_file("service_account.json", scopes=SCOPES)
    # gc = gspread.authorize(creds)
    # sheet = gc.open_by_key(SHEET_ID).sheet1

    # if entry_exists_in_sheet(sheet, palier, seance):
    #     print(f"⏩ Skipping {pptx_path} — already in Google Sheet (Palier {palier}, Séance {seance}).")
    #     return False

    return True

def process_with_ai(pptx_path: str):
    """Send PPTX content to AI and return structured JSON with CSV row + mindmap."""
    print("Processing with AI...")
    title, palier, seance = extract_title_info(pptx_path)
    today = datetime.date.today().strftime("%Y-%m-%d")

    # Load all slide text
    prs = Presentation(pptx_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_content.append(shape.text.strip())
        slides_text.append(f"--- Slide {i} ---\n" + "\n".join(slide_content))
    full_text = "\n\n".join(slides_text)

    # ---------------------------
    # AI PROMPT
    # ---------------------------
    prompt = f"""
You are a teaching assistant generating structured lesson data.

Return **only valid JSON**, no markdown, no extra text.

Follow **exactly** this structure and phrasing style, adapting the content to the slides provided::
{{
  "csv_row": {{
    "Date": "{today}",
    "Parcours": "Parcours 1",
    "Palier": {palier},
    "Séance": {seance},
    "Objectif": "Amener les élèves à comprendre, lire et produire des mots simples du vocabulaire et identifier/écrire les lettres a, i, b, d, e, é, o.",
    "Rituel": "Les élèves chantent « Bonjour les amis », réalisent une dictée flash de mots simples et corrigent collectivement sur ardoise.",
    "Vocabulaire": "Les élèves découvrent et répètent les mots de vocabulaire : ardoise, immeuble, olive, banane, dindon, melon, école, bébé.",
    "Lecture": "Les élèves lisent les lettres a, i, b, d, e, é, o en majuscules et minuscules, puis des mots simples (ananas, immeuble, banane, école).",
    "Écriture": "Les élèves écrivent sur ardoise et dans leur cahier les lettres et mots étudiés, après démonstration au tableau.",
    "Pratique autonome": "Les élèves complètent les activités du livret p.5 et p.7, relient lettres ↔ mots, lisent à voix basse en binômes et recopient les lettres.",
    "Jeu": "Les élèves participent à un jeu « Sauter sur les lettres » : ils sautent sur les cases contenant la lettre annoncée et lisent la syllabe.",
    "Notes": "quelques difficultes dans la séance..."
  }},
  "mindmap": "🧠 **Séance 5 – Palier 2**
🎯 **Objectif principal**
Lire, écrire et utiliser des syllabes et mots simples avec **p – q – qu – t – k – x** à travers du vocabulaire courant.
1. **Rituel**
* Dictée flash de mots connus (sous – sur – lavabo – chat)
* Écriture et correction sur ardoise
2. **Vocabulaire (images + répétition)**
   Mots : poule – coq – taxi – kiwi – savon – se laver – kilo – café
* Répétition chorale et individuelle
* Phrases simples :
  * La poule et le coq.
  * Le taxi est parti.
  * Le garçon se lave les mains avec du savon.
  * Un kilo de tomates.
  * Le café est chaud.
3. **Rebrassage du vocabulaire**
* Jeu « Numéro du mot » : écrire sur ardoise le numéro du mot entendu
* Jeu des images manquantes (taxi – se laver)
* Lecture rapide avec images au tableau
4. **Lecture – Écriture (lettres p, q, qu, t, k, x)**
* Lecture de syllabes : pa – pou – qu – ki – ta – tu – ke – xi…
* Lecture de mots : taxi – kimono – kilo – café – kimono – kiwi
* Lecture de phrases :
  * La dame a acheté un kilo de tomates.
  * Elle a mixé le kiwi.
  * La girafe est avec le petit.
  * Jamal adore le karaté.
* Écriture collective : « Le taxi est parti » (ardoise + cahier)
5. **Pratique autonome**
* Exercices du livret p.15 : activités 1 → 4
* Lecture à voix basse en binômes (échange des rôles)
* Copie de phrases et mots étudiés
6. **Jeu final**
* Jeu du panier des syllabes et lettres (p – q – qu – t – k – x)
* Panier qui circule avec cartes → un élève pioche → lit → la classe répète
7. **Devoir à la maison**
* Écrire les mots du vocabulaire (p.10)
* Relire activités 2 → 4 (p.15)
* Recopier activité 4 (p.15)
* Terminer activité 5 (p.15)
"
}}
Adapt this tone and sentence structure to the actual PPTX content provided below.
Write *original* sentences in the same pedagogical phrasing style — short, action-based, classroom-focused, and aligned with Moroccan French teaching style
Rules:
- Output must be strictly valid JSON that Python's json.loads() can parse.
- Never include ```json or other code fences.
- Escape internal quotes properly.
- Do not add explanations before or after the JSON.
- Base all content on the lesson slides below.

Lesson slides content:
{full_text}
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        temperature=0.1,
        messages=[
            {"role": "system", "content": "You generate structured JSON for a teacher's lesson journal."},
            {"role": "user", "content": prompt}
        ]
    )

    raw_result = response.choices[0].message.content.strip()

    # Try parsing JSON output
    try:
        data = json.loads(raw_result)
        csv_data = data["csv_row"]
        mindmap = data["mindmap"]
    except json.JSONDecodeError as e:
        print("❌ Invalid JSON received:", e)
        print("Raw output:\n", raw_result)
        return None, None

    return csv_data, mindmap

def append_to_google_sheet(json_row: dict):
    """Append a new row to the Google Sheet instead of a CSV file."""
    try:
        SCOPES = ["https://www.googleapis.com/auth/spreadsheets"]
        creds = Credentials.from_service_account_file("service_account.json", scopes=SCOPES)

        # Connect to Sheets
        gc = gspread.authorize(creds)
        sheet = gc.open_by_key(SHEET_ID).sheet1  # Use .worksheet("name") if not the first tab
        existing_data = sheet.get_all_values()

        # If empty → add headers first
        if not existing_data:
            sheet.append_row(HEADERS)
            print("✅ Added headers to new sheet")
        # Ensure all columns are in correct order
        row = [json_row.get(header, "") for header in HEADERS]
        sheet.append_row(row)
        print("✅ Row successfully appended to Google Sheet.")
    except Exception as e:
        print(f"⚠️ Google Sheets sync failed: {e}")
        
    # Define scope for Sheets access

def append_to_csv(csv_data: dict):
    """Append a dict row (12 fields) to CSV safely."""
    print("Adding entry to CSV file...")
    with open(CSV_FILE, "a", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=HEADERS, quoting=csv.QUOTE_ALL)
        writer.writerow(csv_data)

def save_mindmap(pptx_path: str, mindmap: str):
    """Save the mindmap as a Markdown file."""
    title, palier, seance = extract_title_info(pptx_path)
    print("Saving mind map...")
    today = datetime.date.today().strftime("%Y-%m-%d")
    filename = os.path.join(OUTPUT_DIR, f"{title}.md")
    with open(filename, "w", encoding="utf-8") as f:
        f.write(mindmap)
    return filename

# def generate_pdf_from_csv_data(csv_row, pdf_filename):
#     # Load the Jinja2 environment (points to your templates folder)
#     env = Environment(loader=FileSystemLoader("templates"))
#     template = env.get_template("new-template.html")

#     # Render HTML with your data
#     html_content = template.render(csv_row=csv_row)

#     # Make sure output directory exists
#     os.makedirs("output_pdfs", exist_ok=True)
#     pdf_path = os.path.join("output_pdfs", pdf_filename)

#     # Generate PDF
#     HTML(string=html_content).write_pdf(pdf_path)

#     print(f"✅ PDF created: {pdf_path}")

# PROCCES 2 PARCOURS
def process_two_parcours(pptx_path1, pptx_path2):
    csv_data1, mindmap1 = process_with_ai(pptx_path1)
    csv_data2, mindmap2 = process_with_ai(pptx_path2)

    if not csv_data1 or not csv_data2:
        print("❌ Skipped one parcours due to invalid AI output.")
        return

    combined = {
        "Date": csv_data1.get("Date"),
        # parcours 1
        "Palier1": csv_data1.get("Palier"),
        "Séance1": csv_data1.get("Séance"),
        "Objectif1": csv_data1.get("Objectif"),
        "Rituel1": csv_data1.get("Rituel"),
        "Vocabulaire1": csv_data1.get("Vocabulaire"),
        "Lecture1": csv_data1.get("Lecture"),
        "Écriture1": csv_data1.get("Écriture"),
        "Pratique autonome1": csv_data1.get("Pratique autonome"),
        "Jeu1": csv_data1.get("Jeu"),
        "Notes1": csv_data1.get("Notes"),
        # parcours 2
        "Palier2": csv_data2.get("Palier"),
        "Séance2": csv_data2.get("Séance"),
        "Objectif2": csv_data2.get("Objectif"),
        "Rituel2": csv_data2.get("Rituel"),
        "Vocabulaire2": csv_data2.get("Vocabulaire"),
        "Lecture2": csv_data2.get("Lecture"),
        "Écriture2": csv_data2.get("Écriture"),
        "Pratique autonome2": csv_data2.get("Pratique autonome"),
        "Jeu2": csv_data2.get("Jeu"),
        "Notes2": csv_data2.get("Notes"),
    }

    # Save as usual
    append_to_csv(combined)
    # append_to_google_sheet(combined)

    # Combine mindmaps into one markdown file
    mindmap_combined = f"## Parcours 1\n{mindmap1}\n\n---\n\n## Parcours 2\n{mindmap2}"
    title, _, _ = extract_title_info(pptx_path1)
    md_path = os.path.join(OUTPUT_DIR, f"{title}_2parcours.md")
    with open(md_path, "w", encoding="utf-8") as f:
        f.write(mindmap_combined)

    # PDF
    pdf_filename = f"{combined['Palier1']}_P1S{combined['Séance1']}_P2S{combined['Séance2']}.pdf"
    generate_pdf_from_csv_data(combined, pdf_filename)
    print(f"✅ Combined PDF → {pdf_filename}")

def generate_pdf_from_csv_data(csv_row, pdf_filename):
    # 1️⃣ Render HTML with Jinja2
    env = Environment(loader=FileSystemLoader("templates"))
    template = env.get_template("template.html")
    html_content = template.render(csv_row=csv_row)

    # 2️⃣ Save temporary HTML file
    os.makedirs("temp_html", exist_ok=True)
    html_path = os.path.join("temp_html", "temp.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html_content)

    # 3️⃣ Define PDF output path
    os.makedirs("output_pdfs", exist_ok=True)
    pdf_path = os.path.join("output_pdfs", pdf_filename)

    # 4️⃣ Render and export with Playwright (headless Chromium)
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page()
        page.goto(f"file://{os.path.abspath(html_path)}")
        page.pdf(
            path=pdf_path,
            format="A4",
            print_background=True,
            margin={"top": "1cm", "bottom": "1cm", "left": "1cm", "right": "1cm"}
        )
        browser.close()

    print(f"✅ PDF created: {pdf_path}")
    return pdf_path

def choose_pptx_files(pptx_dir):
    """Let the user manually select two PPTX files to pair."""
    all_files = [f for f in os.listdir(pptx_dir) if f.lower().endswith(".pptx")]
    if not all_files:
        print("⚠️ No PPTX files found in directory.")
        exit()

    print("\n📂 Available PPTX files:")
    for i, f in enumerate(all_files, start=1):
        print(f"  {i}. {f}")

    try:
        idx1 = int(input("\nEnter number for Parcours 1 file: ")) - 1
        idx2 = int(input("Enter number for Parcours 2 file: ")) - 1
    except ValueError:
        print("❌ Invalid input. Please enter valid numbers.")
        exit()

    if idx1 not in range(len(all_files)) or idx2 not in range(len(all_files)):
        print("❌ Invalid selection.")
        exit()

    pptx_path1 = os.path.join(pptx_dir, all_files[idx1])
    pptx_path2 = os.path.join(pptx_dir, all_files[idx2])
    print(f"\n✅ Selected:\nParcours 1 → {all_files[idx1]}\nParcours 2 → {all_files[idx2]}")
    return pptx_path1, pptx_path2

# ---------------------------
# MAIN
# ---------------------------
if __name__ == "__main__":
      # or "path/to/folder"
    # pptx_files = [f for f in os.listdir(PPTX_DIR) if f.lower().endswith(".pptx")]
    generate_pdf_from_csv_data(sample_csv_row, "test_layout.pdf")
    # generate_pdf_from_csv_data(csv_data_2parcours, "test_2parcours.pdf")


    ensure_csv_exists()
    if PARCOURS_COUNT == 1:
        pptx_files = [f for f in os.listdir(PPTX_DIR) if f.lower().endswith(".pptx")]
        if not pptx_files:
            print("⚠️ No PPTX files found in directory.")
            exit()
        for pptx_filename in pptx_files:
            pptx_path = os.path.join(PPTX_DIR, pptx_filename)
            # ✅ Skip if already processed (exists in CSV or Google Sheet)
            if not should_process(pptx_path):
                continue
            print("=" * 60)
            print(f"🚀 Processing file: {pptx_filename}")
            print("=" * 60)

            csv_data, mindmap = process_with_ai(pptx_path)

            if csv_data and mindmap:
                append_to_google_sheet(csv_data)
                append_to_csv(csv_data)
                md_file = save_mindmap(pptx_path,mindmap)
                print(f"✅ Processed successfully → {md_file}")
                pdf_filename = f"{csv_data['Palier']}_Séance{csv_data['Séance']}.pdf"
                generate_pdf_from_csv_data(csv_data, pdf_filename)
                print(f"📄 PDF created → {pdf_filename}")

            else:
                print(f"❌ Skipped {pptx_filename} (invalid AI response)")
    else:
         # Manual selection for 2 parcours
        pptx_path1, pptx_path2 = choose_pptx_files(PPTX_DIR)
        process_two_parcours(pptx_path1, pptx_path2)
