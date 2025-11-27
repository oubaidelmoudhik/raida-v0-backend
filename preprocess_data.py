import os
import json
import re
from pptx import Presentation

LESSONS_DIR = "lessons"
OUTPUT_JSON = "data/lessons.json"

def extract_metadata_from_filename(filename: str):
    """
    Extract metadata from filename.
    Supports formats:
    1. FR_N5_P1_SEM1_S3_V2.pptx (Short code)
    2. Français_Niv5_Parcour1_Palier3_Séance1.pptx (Long format)
    """
    name, _ = os.path.splitext(filename)
    parts = name.split("_")
    
    metadata = {
        "title": name.replace("_", " "),
        "subject": "français", # Default
        "level": "5",
        "period": "",
        "week": "",
        "session": ""
    }

    # Try Short Format: FR_N5_P1_SEM1_S3...
    if len(parts) >= 5 and len(parts[0]) <= 4:
        # Subject
        if parts[0].upper() == "FR": metadata["subject"] = "français"
        elif parts[0].upper() == "MATH": metadata["subject"] = "mathématiques"
        elif parts[0].upper() == "AR": metadata["subject"] = "langue arabe"
        
        # Level
        if parts[1].startswith("N"): metadata["level"] = parts[1][1:]
        
        # Period (P)
        if parts[2].startswith("P"): metadata["period"] = parts[2][1:]
        
        # Semaine (SEM) -> Week
        if parts[3].startswith("SEM"): metadata["week"] = parts[3][3:]
        
        # Seance (S) -> Session
        if parts[4].startswith("S"): metadata["session"] = parts[4][1:]

    # Try Long Format: Français_Niv5_Parcour1_Palier3_Séance1
    else:
        for part in parts:
            lower = part.lower()
            if "français" in lower: metadata["subject"] = "français"
            elif "math" in lower: metadata["subject"] = "mathématiques"
            elif "arabe" in lower: metadata["subject"] = "langue arabe"
            
            if lower.startswith("niv"): metadata["level"] = lower.replace("niv", "")
            
            if lower.startswith("parcour"):
                # Map Parcour to Week? Or Period?
                # User example: Parcour1 -> Week 2? 
                # For now, let's map Parcour -> Week as it's the closest equivalent in hierarchy
                metadata["week"] = part.lower().replace("parcour", "").replace("s", "")
            
            if lower.startswith("palier"):
                metadata["period"] = part.lower().replace("palier", "")
                
            if lower.startswith("séance") or lower.startswith("seance"):
                metadata["session"] = part.lower().replace("séance", "").replace("seance", "")

    return metadata

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def extract_objective(content: str):
    """Try to extract objective from content."""
    lines = content.split('\n')
    for line in lines:
        if "objectif" in line.lower() or "capable de" in line.lower():
            # Clean up the line
            return line.strip()
    return "......"

def update_lessons_registry(lessons_dir: str = LESSONS_DIR, json_path: str = OUTPUT_JSON):
    """
    Scan lessons directory for new PPTX files and update lessons.json.
    Returns True if changes were made.
    """
    if not os.path.exists(lessons_dir):
        os.makedirs(lessons_dir, exist_ok=True)
        return False

    # Load existing lessons
    lessons = []
    if os.path.exists(json_path):
        try:
            with open(json_path, "r", encoding="utf-8") as f:
                lessons = json.load(f)
        except json.JSONDecodeError:
            lessons = []
    
    # Get existing filenames
    existing_filenames = {l.get("filename") for l in lessons}
    
    # Scan directory
    files = sorted([f for f in os.listdir(lessons_dir) if f.endswith(".pptx")])
    new_files_found = False

    for filename in files:
        if filename not in existing_filenames:
            print(f"🔍 New lesson found: {filename}")
            pptx_path = os.path.join(lessons_dir, filename)
            
            try:
                content = extract_text_from_pptx(pptx_path)
                meta = extract_metadata_from_filename(filename)
                objective = extract_objective(content)
                
                # Generate new ID
                new_id = max([l["id"] for l in lessons], default=0) + 1
                
                new_lesson = {
                    "id": new_id,
                    "title": meta["title"],
                    "subject": meta["subject"],
                    "level": meta["level"],
                    "period": meta["period"],
                    "week": meta["week"],
                    "session": meta["session"],
                    "filename": filename,
                    "objective": objective,
                    "content": content
                }
                
                lessons.append(new_lesson)
                new_files_found = True
                print(f"✅ Added lesson: {filename}")
            except Exception as e:
                print(f"❌ Error processing {filename}: {e}")

    if new_files_found:
        os.makedirs(os.path.dirname(json_path), exist_ok=True)
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(lessons, f, ensure_ascii=False, indent=2)
        print(f"💾 {json_path} updated")
        return True
    
    return False

def main():
    print("Scanning for new lessons...")
    update_lessons_registry()

if __name__ == "__main__":
    main()
